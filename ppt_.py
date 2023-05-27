import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import MSO_ANCHOR

# Define the path to the input text files and font file
slide1_input_file = 'sample_slide1_input.txt'
slide2_input_file = 'sample_slide2_input.txt'
font_file = 'sample_font_file.ttf'


prs = Presentation()

blank_slide_layout = prs.slide_layouts[6]
slide1 = prs.slides.add_slide(blank_slide_layout)

# Read content from sample_slide1_input.txt
with open(slide1_input_file, 'r') as file:
    slide1_content = file.read()

width = Inches(10)
left = top =  right = Inches(0.1)
txBox = slide1.shapes.add_textbox(left, top, width, right)
tf = txBox.text_frame

tf.text = slide1_content
tf.word_wrap = True

# Set the font style for Slide 1
for paragraph in tf.paragraphs:
    for run in paragraph.runs:
        run.font.file = font_file
        run.font.name = 'Love Ya Like A Sister'
        run.font.size = Pt(13)
        
# p = tf.add_paragraph()
# p.text = "This is a second paragraph that's bold"
# p.font.bold = True

# p = tf.add_paragraph()
# p.text = "This is a third paragraph that's big"
# p.font.size = Pt(40)

blank_slide_layout = prs.slide_layouts[6]
slide2 = prs.slides.add_slide(blank_slide_layout)

# Read content from sample_slide1_input.txt
with open(slide2_input_file, 'r') as file:
    slide2_content = file.read()

width = Inches(10)
left = top =  right  = Inches(0.1)
txBox = slide2.shapes.add_textbox(left, top, width, right)
tf = txBox.text_frame

tf.text = slide2_content
tf.word_wrap = True

# Set the font style for Slide 1
for paragraph in tf.paragraphs:
    for run in paragraph.runs:
        run.font.file = font_file
        run.font.name = 'Love Ya Like A Sister'
        run.font.size = Pt(23)


prs.save('test.pptx')
